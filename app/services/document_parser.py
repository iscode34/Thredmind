import io
from pathlib import Path



def extract_text_from_file(file_bytes: bytes, filename: str) -> tuple[str, str]:
    ext = Path(filename).suffix.lower()
    if ext == ".pdf":
        return _extract_pdf(file_bytes), "pdf"
    elif ext == ".docx":
        return _extract_docx(file_bytes), "docx"
    elif ext == ".pptx":
        return _extract_pptx(file_bytes), "pptx"
    elif ext in (".txt", ".md"):
        return file_bytes.decode("utf-8", errors="replace"), ext.lstrip(".")
    else:
        raise ValueError(f"Unsupported file format: {ext}")


def _extract_pdf(file_bytes: bytes) -> str:
    import PyPDF2
    reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
    text_parts = []
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text_parts.append(page_text)
    return "\n\n".join(text_parts)


def _extract_docx(file_bytes: bytes) -> str:
    from docx import Document as DocxDocument
    doc = DocxDocument(io.BytesIO(file_bytes))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_pptx(file_bytes: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))

    text_parts = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        slide_text.append(para.text.strip())
        if slide_text:
            text_parts.append("\n".join(slide_text))
    return "\n\n".join(text_parts)
